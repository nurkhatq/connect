import os
import re
import asyncio
import aiofiles
import numpy as np
from pathlib import Path
from typing import List, Dict, Any
from sentence_transformers import SentenceTransformer
import logging

logger = logging.getLogger(__name__)

# Глобальная модель для embeddings
_sentence_model = None

def get_sentence_model():
    """Ленивая загрузка модели"""
    global _sentence_model
    if _sentence_model is None:
        _sentence_model = SentenceTransformer('all-MiniLM-L12-v2')
    return _sentence_model

def extract_sources_list(source_docs) -> List[str]:
    """Извлекает список уникальных источников"""
    seen = set()
    sources = []
    for doc in source_docs:
        file_name = doc.metadata.get("file_name")
        if file_name and file_name not in seen:
            seen.add(file_name)
            sources.append(file_name)
    return sources

async def extract_text_from_file_async(filepath: str) -> str:
    """Асинхронное извлечение текста из файла"""
    ext = os.path.splitext(filepath)[1].lower()
    
    if ext == '.txt':
        # Простой текстовый файл читаем асинхронно
        import chardet
        
        # Определяем кодировку
        async with aiofiles.open(filepath, 'rb') as f:
            raw_data = await f.read(10000)
        
        detected = chardet.detect(raw_data)
        encoding = detected['encoding'] or 'utf-8'
        
        async with aiofiles.open(filepath, 'r', encoding=encoding, errors='ignore') as f:
            return await f.read()
    
    else:
        # Для остальных форматов используем синхронную версию в executor
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, extract_text_from_file, filepath)

def extract_text_from_file(filepath: str) -> str:
    """Синхронное извлечение текста (для executor)"""
    ext = os.path.splitext(filepath)[1].lower()
    
    try:
        if ext == '.txt':
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
                
        elif ext == '.docx':
            from docx import Document
            doc = Document(filepath)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            tables = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        tables.append(" | ".join(row_text))
            return '\n'.join(paragraphs + tables)
            
        elif ext == '.pdf':
            import PyPDF2
            text = ''
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() or ''
            return text
            
        elif ext in ['.xlsx', '.xls']:
            import openpyxl
            wb = openpyxl.load_workbook(filepath, data_only=True)
            text = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) for cell in row if cell is not None]
                    if row_text:
                        text.append(' | '.join(row_text))
            return '\n'.join(text)
            
        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(filepath)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        t = shape.text.strip()
                        if t:
                            text.append(t)
            return '\n'.join(text)
            
    except Exception as e:
        logger.error(f"Error extracting text from {filepath}: {e}")
        return ''
    
    return ''

def cosine_similarity(a: np.ndarray, b: np.ndarray) -> float:
    """Вычисляет косинусное сходство"""
    if not a.any() or not b.any():
        return 0.0
    return float(np.dot(a, b) / (np.linalg.norm(a) * np.linalg.norm(b)))

def normalize_text(text: str) -> str:
    """Нормализует текст"""
    text = text.lower().strip()
    text = re.sub(r'\s+', ' ', text)
    return text

async def find_similar_files_async(uploaded_text: str, folder: str, threshold: float = 0.7) -> List[Dict]:
    """Асинхронный поиск похожих файлов"""
    model = get_sentence_model()
    uploaded_text_norm = normalize_text(uploaded_text)
    
    # Получаем embedding в executor
    loop = asyncio.get_event_loop()
    uploaded_emb = await loop.run_in_executor(
        None,
        lambda: model.encode([uploaded_text_norm], convert_to_numpy=True)[0]
    )
    
    similar = []
    folder_path = Path(folder)
    
    # Собираем все файлы
    files = [f for f in folder_path.glob("**/*") 
             if f.is_file() and f.suffix.lower() in [".docx", ".pdf", ".txt", ".xlsx", ".pptx"]]
    
    # Обрабатываем файлы параллельно
    tasks = []
    for file_path in files:
        tasks.append(process_file_similarity(file_path, model, uploaded_emb, threshold))
    
    results = await asyncio.gather(*tasks, return_exceptions=True)
    
    for result in results:
        if isinstance(result, dict) and result['similarity'] >= threshold:
            similar.append(result)
    
    # Сортируем по схожести
    similar.sort(key=lambda x: x['similarity'], reverse=True)
    
    # Возвращаем точные совпадения или топ-3
    exact_matches = [f for f in similar if f['similarity'] == 100.0]
    if exact_matches:
        return exact_matches
    
    return similar[:3]

async def process_file_similarity(file_path: Path, model, uploaded_emb: np.ndarray, threshold: float) -> Dict:
    """Обрабатывает один файл для проверки схожести"""
    try:
        # Извлекаем текст
        text = await extract_text_from_file_async(str(file_path))
        if not text.strip():
            return {'file': file_path.name, 'similarity': 0}
        
        text_norm = normalize_text(text)
        
        # Получаем embedding
        loop = asyncio.get_event_loop()
        emb = await loop.run_in_executor(
            None,
            lambda: model.encode([text_norm], convert_to_numpy=True)[0]
        )
        
        # Вычисляем схожесть
        sim = cosine_similarity(uploaded_emb, emb)
        
        return {
            'file': file_path.name,
            'similarity': round(sim * 100, 2)
        }
        
    except Exception as e:
        logger.error(f"Error processing {file_path}: {e}")
        return {'file': file_path.name, 'similarity': 0}